"""Builds the 2026-08-22 semantic-guidance status deck (14 slides, dark theme).

Audience: two professors NOT familiar with the experiment history. The slides stay sparse;
the speaker notes carry the teaching. Notes follow a strict one-bullet-per-line rule so they
can be scanned while presenting - never wrap a bullet.

Styling copied from reports/presentations/2026-05_E3a-status.pptx (10 x 5.62 in, #1C2340
background, cyan/navy palette).

All confusion-matrix numbers are asserted against the per-clip result files before the deck
is written, so a stale figure can never be silently embedded.

    python student_training/scripts/build_status_presentation_2026-08-22.py
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT / "student_training" / "scripts"))
from metrics_core import metrics_from_arrays  # noqa: E402

FIG = ROOT / "reports" / "figures"
OUT = ROOT / "reports" / "presentations" / "2026-08_semantic-guidance-status.pptx"

# ---- E3a-status palette
BG = RGBColor(0x1C, 0x23, 0x40)
PANEL = RGBColor(0x24, 0x30, 0x60)
PANEL_DK = RGBColor(0x20, 0x2A, 0x50)
CYAN = RGBColor(0x00, 0xBF, 0xFF)
CYAN_DK = RGBColor(0x00, 0x96, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xA0, 0xB4, 0xCC)
GREEN = RGBColor(0x00, 0xE6, 0x76)
ORANGE = RGBColor(0xFF, 0xA7, 0x26)
RED = RGBColor(0xFF, 0x6B, 0x6B)
SW, SH = 10.0, 5.625

FOOTER = "CCP-MMLM · Semantic Guidance for Collision Anticipation"

# ---------------------------------------------------------------- verified data
# arm, results file, TP, FN, FP, TN, P, R, F1, Acc   (all at threshold 0.5)
CM_ROWS = [
    ("A0 · Baseline", "StageA_scorer/badas_open_private.jsonl",
     308, 30, 130, 209, .703, .911, .794, .764),
    ("A1 · Crash-only (control)", "a1_1761/test_results_ep04.jsonl",
     320, 18, 123, 216, .722, .947, .820, .792),
    ("B-v1 · Joint, leaky captions", "b_1761_par/test_results_ep04.jsonl",
     317, 21, 130, 209, .709, .938, .808, .777),
    ("B-v2 · Joint, clean captions", "b_v2_1761/test_results_ep02.jsonl",
     285, 53, 76, 263, .789, .843, .816, .809),
    ("B-v3 · Joint, fixed procedure", "b_v3_1761/test_results_ep10.jsonl",
     267, 71, 55, 284, .829, .790, .809, .814),
    ("P1 · Two-stage", "p1_stageB/test_results_ep02.jsonl",
     278, 60, 92, 247, .751, .823, .785, .775),
]

# Arm names MUST match CM_ROWS exactly - the two tables are read side by side.
EXPERIMENTS = [
    ("A0 · Baseline",
     "Public BADAS-Open checkpoint, no training by us. The bar to beat.", "0.853", "0.864"),
    ("A1 · Crash-only (control)",
     "LoRA fine-tune on collision labels alone. No language anywhere.", "0.900", "0.904"),
    ("B-v1 · Joint, leaky captions",
     "First semantic arm. Captions were later shown to encode the label.", "0.890", "0.895"),
    ("B-v2 · Joint, clean captions",
     "Same design on de-leaked captions. The gap grew instead of closing.", "0.880", "0.891"),
    ("B-v3 · Joint, fixed procedure",
     "Predictor warm-started as the protocol specified + per-group gradient "
     "clipping so neither loss can dominate a step.", "0.877", "0.887"),
    ("P1 · Two-stage",
     "Language first to convergence, then crash only. Worst of all.", "0.827", "0.848"),
]


def verify():
    """Fail loudly rather than embed a stale number."""
    base = ROOT / "outputs" / "e4_vjepa_reason"
    for name, rel, tp, fn, fp, tn, p, r, f1, acc in CM_ROWS:
        rows = [json.loads(l) for l in open(base / rel, encoding="utf-8") if l.strip()]
        m = metrics_from_arrays([int(x["ground_truth"]) for x in rows],
                                [float(x["score"]) for x in rows], threshold=0.5)
        got = (m["tp"], m["fn"], m["fp"], m["tn"])
        if got != (tp, fn, fp, tn):
            raise SystemExit(f"CM MISMATCH {name}: file {got}, deck {(tp, fn, fp, tn)}")
        for label, want, have in (("P", p, m["precision"]),
                                  ("R", r, m["recall_sensitivity_tpr"]),
                                  ("F1", f1, m["f1"]), ("Acc", acc, m["accuracy"])):
            if abs(want - have) > 0.0006:
                raise SystemExit(f"{label} MISMATCH {name}: file {have:.4f} vs deck {want}")
    names_cm = [r[0] for r in CM_ROWS]
    names_ex = [r[0] for r in EXPERIMENTS]
    if names_cm != names_ex:
        raise SystemExit(f"ARM NAMES DIVERGE:\n  {names_cm}\n  {names_ex}")
    print("  [verify] CM + P/R/F1/Acc reproduce from per-clip files; arm names consistent")


# ---------------------------------------------------------------- primitives
def _bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def rect(slide, x, y, w, h, color):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def text(slide, x, y, w, h, runs, size=11, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
         font="Calibri", italic=False, space_after=4, line_spacing=None):
    """runs: str, or list of paragraphs; each paragraph is a str or a list of (txt, dict)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    paras = [runs] if isinstance(runs, str) else runs
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        pieces = [(para, {})] if isinstance(para, str) else para
        for txt, ov in pieces:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(ov.get("size", size))
            r.font.bold = ov.get("bold", bold)
            r.font.italic = ov.get("italic", italic)
            r.font.color.rgb = ov.get("color", color)
            r.font.name = ov.get("font", font)
    return tb


def notes(slide, lines):
    """One bullet per line. Never wraps - the user reads these live."""
    tf = slide.notes_slide.notes_text_frame
    tf.text = lines[0]
    for l in lines[1:]:
        tf.add_paragraph().text = l


def picture(slide, path, x, y, w):
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))


def content_slide(prs, title, subline=None, page=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    rect(s, 0.35, 0.22, 0.07, 0.50, CYAN)
    # 25pt fits ~52 chars on a 9.2in box; shrink rather than wrap into the rule below
    t_size = 25 if len(title) <= 52 else max(19, 25 * 52 / len(title))
    text(s, 0.50, 0.13, 9.20, 0.60, title, size=t_size, color=CYAN, bold=True)
    rect(s, 0.35, 0.84, 9.30, 0.012, CYAN_DK)
    if subline:
        text(s, 0.50, 0.88, 9.20, 0.40, subline, size=11, color=MUTED)
    text(s, 0.50, 5.29, 7.00, 0.30, FOOTER, size=8.5, color=MUTED)
    if page:
        text(s, 8.80, 5.29, 0.80, 0.30, str(page), size=8.5, color=MUTED,
             align=PP_ALIGN.RIGHT)
    return s


def card(slide, x, y, w, h, title, bullets, accent=CYAN, sub=None, body_size=9.5):
    rect(slide, x, y, w, h, PANEL)
    rect(slide, x, y, w, 0.05, accent)
    text(slide, x + 0.13, y + 0.10, w - 0.26, 0.30, title, size=12, color=accent, bold=True)
    top = y + 0.42
    if sub:
        text(slide, x + 0.13, y + 0.38, w - 0.26, 0.26, sub, size=9, color=MUTED)
        top = y + 0.68
    text(slide, x + 0.13, top, w - 0.26, h - (top - y) - 0.08,
         [f"• {b}" for b in bullets], size=body_size, color=WHITE, space_after=3)


def hero(slide, x, y, w, h, headline, tail, accent=GREEN):
    rect(slide, x, y, w, h, PANEL_DK)
    rect(slide, x, y, 0.07, h, accent)
    text(slide, x + 0.20, y + 0.08, w - 0.30, h - 0.12,
         [[(headline, {"size": 21, "bold": True, "color": accent}),
           (tail, {"size": 12, "color": WHITE})]], size=12)


def table(slide, rows, cols_w, x, y, h, font=8.5, header_font=8.5, align_first_left=True):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y),
                                   Inches(sum(cols_w)), Inches(h))
    tbl = shape.table
    tbl.first_row = False
    for j, w in enumerate(cols_w):
        tbl.columns[j].width = Inches(w)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            # python-pptx has no API for cell fill colour -> drop in the XML directly
            fill = "202A50" if i == 0 else ("243060" if i % 2 else "1F2950")
            cell._tc.get_or_add_tcPr().append(
                parse_xml(f'<a:solidFill {nsdecls("a")}><a:srgbClr val="{fill}"/></a:solidFill>'))
            cell.text = str(val)
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            for p in cell.text_frame.paragraphs:
                p.alignment = (PP_ALIGN.LEFT if (j == 0 and align_first_left)
                               else PP_ALIGN.CENTER)
                for r in p.runs:
                    r.font.size = Pt(header_font if i == 0 else font)
                    r.font.name = "Calibri"
                    r.font.bold = (i == 0)
                    r.font.color.rgb = CYAN if i == 0 else WHITE
    return tbl


# ================================================================= SLIDES
def s01_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    rect(s, 0, 0, SW, 0.08, CYAN)
    rect(s, 0, SH - 0.08, SW, 0.08, CYAN)
    rect(s, 0, 0, 0.08, SH, CYAN_DK)
    rect(s, 7.40, 0, 2.60, 2.60, PANEL_DK)
    text(s, 0.50, 0.80, 9.00, 1.10, "CCP-MMLM", size=48, color=WHITE, bold=True)
    rect(s, 0.50, 1.90, 6.50, 0.06, CYAN)
    text(s, 0.50, 1.99, 8.60, 0.90,
         "Semantic Guidance for Collision Anticipation\nExplainable dashcam collision "
         "prediction via Teacher-Student distillation", size=15, color=MUTED)
    rect(s, 0.50, 3.05, 9.00, 1.42, PANEL)
    rect(s, 0.50, 3.05, 9.00, 0.05, CYAN)
    text(s, 0.62, 3.13, 8.70, 0.30, "Research goal", size=12, color=CYAN, bold=True)
    text(s, 0.62, 3.45, 8.76, 1.00,
         "Teach a compact vision model to anticipate collisions by supervising it, during "
         "training only, with sentences a large vision-language model writes about each clip. "
         "A yes/no label carries about one bit; a sentence carries far more. The language "
         "machinery is removed after training, so the deployed model is vision-only - same "
         "cost, same latency.", size=11, color=WHITE)
    text(s, 0.50, 4.62, 7.00, 0.35, "MSc Thesis - Status Update  |  August 2026",
         size=11.5, color=CYAN)
    text(s, 0.50, 4.92, 5.00, 0.42, "Eviatar Ohayon", size=15, color=WHITE, bold=True)
    notes(s, [
        "Opening framing:",
        "- Task: predict from dashcam video that a collision is about to happen.",
        "- Idea: during training, also make the model reproduce the meaning of a teacher's sentence.",
        "- Why: a yes/no label is ~1 bit per example; a sentence describing the scene is far richer.",
        "- Key property: the language parts are deleted after training - zero runtime cost.",
        "- Formal name for this setting: learning using privileged information (LUPI).",
        "- Today: where it stands after six arms, plus what the next two weeks target."])


def s02_outcome(prs):
    s = content_slide(prs, "Outcome - Where Semantic Guidance Stands",
                      "One banked win from vision alone; language demonstrably learns, but has "
                      "not yet moved the collision metric.", page=2)
    hero(s, 0.50, 1.28, 9.00, 0.72, "Test AP = 0.900",
         "   on 677 held-out clips   -   +0.047 over the public baseline (0.853)")
    card(s, 0.50, 2.22, 2.93, 2.10, "[DONE]  Control",
         ["Crash labels only, no language", "AP 0.900 / AUC 0.904",
          "Beats published BADAS-Open", "This result is banked"],
         accent=GREEN, sub="A1 - crash-only LoRA")
    card(s, 3.54, 2.22, 2.93, 2.10, "[OPEN]  Semantic arms",
         ["Four joint arms + one two-stage", "Best 0.890, worst 0.827",
          "All lose to the control", "Gap widened as bugs were fixed"],
         accent=ORANGE, sub="B-v1 / v2 / v3 - P1")
    card(s, 6.57, 2.22, 2.93, 2.10, "[RESULT]  Representation",
         ["Frozen features encode captions", "31x chance on unseen clips",
          "Control sits exactly at chance", "Still scaling with more data"],
         accent=CYAN, sub="frozen-trunk probe")
    notes(s, [
        "Three things, and they answer different questions:",
        "- LEFT: our fine-tuned model reaches AP 0.900 on 677 fully held-out clips. Banked.",
        "- It beats the published BADAS-Open baseline (0.853) by +0.047. This arm uses NO language.",
        "- MIDDLE: every semantic arm lost. Best 0.890, worst 0.827 - below even the baseline.",
        "- The gap got WIDER each time we fixed a defect. That pattern is itself the finding.",
        "- RIGHT: a probe shows the frozen features already encode caption meaning at 31x chance.",
        "- Be direct: that is a representation result, not yet a collision-accuracy result.",
        "- AP = area under the precision-recall curve. Threshold-free, so no alarm level is assumed."])


def s03_idea(prs):
    s = content_slide(prs, "The Hypothesis - Language as a Training-Only Teacher",
                      "A sentence is a far richer training signal than a label - but it is "
                      "scaffolding, removed before deployment.", page=3)
    picture(s, FIG / "arch_L1_overview_2026-08-22.png", 0.42, 1.28, 9.16)
    notes(s, [
        "The argument for the method:",
        "- A yes/no collision label carries about one bit of information per training example.",
        "- A sentence naming the vehicle, its motion and whether the gap is closing carries far more.",
        "- So during training the model does two jobs: predict the collision AND match the sentence.",
        "- The sentence comes from a large vision-language model acting as teacher.",
        "- After training the Predictor and text encoder are deleted - the model is vision-only.",
        "- No text, no language model, no extra latency and no extra dependency at run time.",
        "- Expected mechanism: language should shape features to encode WHY a scene is dangerous."])


def s04_inference(prs):
    s = content_slide(prs, "Inference Architecture - Vision-Only, Zero Added Cost",
                      "What actually ships: the BADAS-Open stack with our LoRA folded into the "
                      "base weights.", page=4)
    picture(s, FIG / "arch_L2_inference_2026-08-22.png", 0.42, 1.35, 9.16)
    notes(s, [
        "Walk the blocks left to right:",
        "- 16 frames at 1280x720 are squash-resized to 256x320 and ImageNet-normalised.",
        "- V-JEPA2 ViT-L: 24-layer video transformer, 334M params, frozen; LoRA adds 2.8M (0.84%).",
        "- It emits a grid of 2,560 spatiotemporal patch tokens, each 1,024-dimensional.",
        "- Temporal processor: an attentive probe - ONE learned query attends over all 2,560 tokens.",
        "- That is a 2560-to-1 compression down to a single 1,024-d vector. Frozen.",
        "- Classifier: MLPHead, 3 layers - [Linear-GELU-LayerNorm-Dropout] x2 then Linear. Frozen.",
        "- Output is 2 logits; softmax gives P(collision), a single scalar risk score.",
        "- LoRA MERGES into the base weights (W becomes W + BA) after training.",
        "- So the shipped model is architecturally identical to stock BADAS-Open - only values differ.",
        "- Same parameter count, same latency, same dependencies as the public checkpoint."])


def s05_training(prs):
    s = content_slide(prs, "Training Architecture - One Trunk, Two Objectives",
                      "The same vision path, plus a second branch that exists only while "
                      "training.", page=5)
    picture(s, FIG / "arch_L3_training_2026-08-22.png", 0.42, 1.20, 9.16)
    notes(s, [
        "Only what is new versus the previous slide:",
        "- The vision path is unchanged - same trunk, same frozen pooling and classifier.",
        "- A second branch taps the SAME 2,560-token patch grid, before any pooling.",
        "- Predictor: a small resampler, 1.25M params, with 8 learned queries. TRAINABLE.",
        "- It outputs (1, 8, 768), mean-pooled over the queries to a single 768-d vector.",
        "- The teacher caption goes through a frozen SigLIP text encoder to a 768-d target.",
        "- L_sem compares those two vectors; L_crash compares P(collision) against the label.",
        "- Trainable in total: LoRA adapters (2.8M) + Predictor (1.25M). Everything else frozen.",
        "- The classifier NEVER moves - it stays fitted to the pre-LoRA encoder. See slide 13.",
        "- The whole shaded region is discarded at inference."])


def s06_objective(prs):
    s = content_slide(prs, "The Objective - Crash Loss + Contrastive Semantic Loss",
                      "One weighted sum. The control arm is the same equation with the second "
                      "term switched off.", page=6)
    rect(s, 0.50, 1.30, 9.00, 0.80, PANEL_DK)
    rect(s, 0.50, 1.30, 0.07, 0.80, ORANGE)
    text(s, 0.75, 1.46, 8.60, 0.50, "L  =  L_crash  +  0.05 x L_sem",
         size=19, color=ORANGE, bold=True, font="Consolas")
    text(s, 0.75, 1.83, 8.60, 0.25, "the control arm (A1) sets the second term to zero",
         size=9.5, color=MUTED, italic=True)

    rect(s, 0.50, 2.28, 4.38, 2.00, PANEL)
    rect(s, 0.50, 2.28, 4.38, 0.05, CYAN)
    text(s, 0.68, 2.38, 4.05, 0.30, "L_crash   -   cross-entropy", size=12, color=CYAN, bold=True)
    text(s, 0.68, 2.76, 4.10, 0.45, "CE(z, y)  =  - log softmax(z)_y",
         size=12, color=WHITE, font="Consolas")
    text(s, 0.68, 3.24, 4.10, 0.60,
         "z = the 2 logits, y = the ground-truth label.\nAsks: is this clip a collision?",
         size=9.5, color=MUTED)

    rect(s, 5.12, 2.28, 4.38, 2.00, PANEL)
    rect(s, 5.12, 2.28, 4.38, 0.05, ORANGE)
    text(s, 5.30, 2.38, 4.05, 0.30, "L_sem   -   InfoNCE", size=12, color=ORANGE, bold=True)
    text(s, 5.30, 2.76, 4.15, 0.45, "- log [ exp(p.t+/T) / SUM_j exp(p.t_j/T) ]",
         size=10.5, color=WHITE, font="Consolas")
    text(s, 5.30, 3.24, 4.15, 0.60,
         "p = predicted vector, t = caption vectors, T = learned temperature.\n"
         "Asks: is p closest to THIS clip's caption?", size=9.5, color=MUTED)

    text(s, 0.50, 4.48, 9.00, 0.55,
         "Both p and t are L2-normalised, so p.t is a cosine. InfoNCE scores a candidate not by "
         "how close it is, but by how much closer it is than every alternative.",
         size=10.5, color=WHITE)
    notes(s, [
        "Two questions to expect, both answered here:",
        "WHY CONTRASTIVE AND NOT PLAIN SIMILARITY:",
        "- The first version simply asked the output to be CLOSE to the caption embedding.",
        "- That has a shortcut: always emit the AVERAGE caption. Caption embeddings cluster tightly.",
        "- Measured: the trained model beat that do-nothing baseline by 0.53% of the available range.",
        "- Essentially zero learning, hidden behind a healthy-looking loss curve.",
        "- InfoNCE removes the shortcut by construction, and here is the actual mechanism:",
        "- The mean-caption direction is SHARED by every candidate in the softmax denominator.",
        "- A shared component contributes equally to every term, so it CANCELS in the ratio.",
        "- The collapse solution therefore scores at chance under InfoNCE instead of at 0.865.",
        "WHERE 0.05 COMES FROM:",
        "- The two losses have very different natural magnitudes, so they must be rescaled.",
        "- Cosine loss sits around 0.13; InfoNCE starts near log(N) = 7.25 nats at N = 1,413.",
        "- 0.05 x 7.25 is about 0.35, roughly a third of the crash loss (about 1.0 at init).",
        "- It is a magnitude-matching heuristic - it was NEVER swept. A fair reviewer objection.",
        "- InfoNCE = Information Noise-Contrastive Estimation."])


def s07_data(prs):
    s = content_slide(prs, "Data - 4,446 Windows and How the Captions Were Written",
                      "The vision model trains on all windows; only a curated 1,761 carry "
                      "teacher captions.", page=7)
    picture(s, FIG / "dataset_pipeline_2026-08-22.png", 0.42, 1.18, 9.16)
    rect(s, 0.50, 3.72, 9.00, 1.42, PANEL)
    rect(s, 0.50, 3.72, 9.00, 0.05, ORANGE)
    text(s, 0.68, 3.80, 8.70, 0.28, "Same clip, old caption vs new caption",
         size=11, color=ORANGE, bold=True)
    text(s, 0.68, 4.10, 8.75, 0.45,
         [[("V10 (leaky):  ", {"size": 9.5, "bold": True, "color": RED}),
           ("\"Ego vehicle closing distance to a stationary silver sedan parked on the right "
            "side of a narrow urban street while driving forward.\"",
            {"size": 9.5, "color": WHITE})]])
    text(s, 0.68, 4.58, 8.75, 0.45,
         [[("V12 (neutral):  ", {"size": 9.5, "bold": True, "color": GREEN}),
           ("\"Gray hatchback driving forward directly ahead in ego lane maintains a constant "
            "following distance.\"", {"size": 9.5, "color": WHITE})]])
    notes(s, [
        "How the caption set was actually built:",
        "- The full window pool is 1,482 clips x 3 windows = 4,446, exactly balanced YES/NO.",
        "- The vision-only arms train on windows; captions are needed only for the semantic arms.",
        "- We ran the frozen baseline over all 4,446 and mined the 587 windows it got WRONG.",
        "- Captioning only failures would bias the set, so we added 1,174 windows it got RIGHT.",
        "- That gives 1,761 windows to caption, at 2 parts correct to 1 part mined failure.",
        "- Teacher: google/gemini-3.6-flash via OpenRouter. 1,761 captions cost $32.82.",
        "- It won a prompt bake-off: 72.2% accuracy, zero false positives, best caption fidelity.",
        "- Captions are stored as JSONL, one row per window, embedded once by frozen SigLIP.",
        "ON THE TWO EXAMPLES:",
        "- Both rows describe the SAME window, and both clips end in a collision.",
        "- V10 says 'closing distance' - outcome language. A text-only classifier reached AUC 0.964.",
        "- V12 only describes what is visible, so the caption cannot stand in for the label.",
        "- That matters: our hypothesis needs the caption to ADD information, not repeat the label.",
        "ON BALANCE, IF ASKED:",
        "- Positives use 0.5 / 1.0 / 1.5 s before the event: 253 / 268 / 335 - roughly, not exactly.",
        "- Negatives have no event to count down to, so they use midpoint offsets: 287 / 298 / 320."])


def s08_experiments(prs):
    s = content_slide(prs, "Experiment Log - Six Arms, One Control",
                      "Each arm fixed the previous arm's defect. The gap to the control widened "
                      "every time.", page=8)
    rows = [["Arm", "What it tested", "test AP", "AUC"]]
    for name, desc, ap, auc in EXPERIMENTS:
        rows.append([name, desc, ap, auc])
    table(s, rows, [2.05, 5.35, 0.80, 0.80], 0.50, 1.32, 3.25, font=8.2, header_font=8.5)
    text(s, 0.50, 4.78, 9.00, 0.30,
         "All scored on the same 677 held-out clips. Paired bootstrap, 5,000 resamples: the "
         "control stayed ahead in 100% of draws in every comparison.", size=9, color=MUTED)
    notes(s, [
        "Read it as a story, not a list:",
        "- A0 is the published baseline, untrained by us: AP 0.853.",
        "- A1 is our crash-only fine-tune: 0.900. Every semantic arm has to beat this.",
        "- Then four semantic attempts, each fixing what the previous got wrong. All lost.",
        "- The gap grew: 0.010, then 0.019, then 0.022, then 0.073 for the two-stage arm.",
        "WHAT V10 WAS AND WHY IT LEAKED:",
        "- V10 used two prompt modes: GT-informed for positives, blind for negatives.",
        "- That produced two different vocabularies by class, so the wording encoded the answer.",
        "- A plain TF-IDF classifier recovered the label from caption TEXT ALONE at AUC 0.964.",
        "- Higher than the vision model itself - the semantic loss was a noisy copy of the label.",
        "WHAT V12 CHANGED:",
        "- One neutral prompt for both classes, with no ground-truth branch at all.",
        "- A closed four-word vocabulary for gap dynamics, plus symmetric alarm/outcome word bans.",
        "- Leakage dropped 0.964 to 0.764. Target was below 0.75, so a narrow, accepted miss.",
        "- Crucially, cleaning the captions made the AP gap WIDER - so leakage was not the cause.",
        "WHAT B-v3 FIXED:",
        "- Our protocol said to warm-start the Predictor from the probe; that step had been skipped.",
        "- Added per-group gradient clipping so neither loss can dominate a single step.",
        "- Both defects were real, both were corrected - and the arm still lost."])


def s09_curves(prs):
    s = content_slide(prs, "Training Dynamics - Every Arm Overfits by Epoch 2",
                      "Validation loss bottoms out early in all four arms; checkpoints are "
                      "selected on validation AP, not loss.", page=9)
    picture(s, FIG / "loss_curves" / "loss_grid_4arms_2026-08-22.png", 1.15, 1.20, 7.70)
    notes(s, [
        "What to point at:",
        "- Same shape everywhere: train loss falls, val loss bottoms near epoch 2, then climbs.",
        "- That is overfitting - the model is memorising 1,761 windows.",
        "- The orange line is the checkpoint we keep, chosen on validation AP rather than loss.",
        "- Loss and AP disagree: AP holds up for a couple of epochs after loss has already turned.",
        "- Loss punishes overconfidence everywhere; AP only cares about the ranking.",
        "- Bottom-right is the two-stage arm - its train/val gap reaches double the control's.",
        "- Cause: warm-started from an adapted checkpoint but reused the from-scratch learning rate.",
        "- That is a correctable procedural choice, not a verdict on two-stage training."])


def s10_metrics(prs):
    s = content_slide(prs, "Performance Metrics - Confusion Matrices at Threshold 0.5",
                      "Same arms as the experiment log, same 677 clips, one uniform threshold.",
                      page=10)
    rows = [["Arm", "TP", "FN", "FP", "TN", "Prec", "Rec", "F1", "Acc"]]
    for name, _, tp, fn, fp, tn, p, r, f1, acc in CM_ROWS:
        rows.append([name, tp, fn, fp, tn, f"{p:.3f}", f"{r:.3f}", f"{f1:.3f}", f"{acc:.3f}"])
    table(s, rows, [2.30, 0.55, 0.55, 0.55, 0.55, 0.75, 0.72, 0.72, 0.72], 0.50, 1.32, 2.95,
          font=8.5, header_font=8.5)
    text(s, 0.50, 4.42, 9.00, 0.30, "677 clips - 338 collision / 339 safe.",
         size=9, color=MUTED)
    text(s, 0.50, 4.70, 9.00, 0.55,
         "The semantic arms look more precise here only because their scores sit lower - at "
         "MATCHED precision the control still wins on recall. AP and AUC are the honest comparison.",
         size=9.5, color=RED, italic=True)
    notes(s, [
        "Expect this question, and it is a good one:",
        "- 'B-v3 has fewer total errors (126) and a nicer FP/FN split than A1 (141). How is it worse?'",
        "- Because 0.5 is arbitrary: it lands each model at a different point on its own PR curve.",
        "- A1's scores sit HIGH, so at 0.5 it over-calls collisions: 123 false positives.",
        "- That is a calibration offset, not a ranking failure - and you fix it by moving the threshold.",
        "- At each arm's OWN best threshold the ordering flips straight back:",
        "-     A1   : 141 errors at 0.5  ->  115 errors at 0.647   (accuracy 0.792 -> 0.830)",
        "-     B-v2 : 129 errors at 0.5  ->  128 errors at 0.491   (accuracy 0.809 -> 0.811)",
        "-     B-v3 : 126 errors at 0.5  ->  117 errors at 0.297   (accuracy 0.814 -> 0.827)",
        "-     P1   : 152 errors at 0.5  ->  140 errors at 0.738   (accuracy 0.775 -> 0.793)",
        "- A1 ends up with FEWER errors than B-v3 ever achieves: 115 against 117.",
        "- Direct check: force A1 to B-v3's precision (0.829) and it still wins recall, 0.805 vs 0.790.",
        "- Analogy: two graders, one generous one strict. At a fixed pass mark the strict looks better.",
        "- That says nothing about which grader RANKS students correctly. AP measures the ranking."])


def s11_probe(prs):
    s = content_slide(prs, "Representation Probe - What Frozen Features Encode",
                      "A frozen-trunk diagnostic: no LoRA, no crash training. It measures what "
                      "BADAS already contains.", page=11)
    picture(s, FIG / "semantic_retrieval_scaling_2026-08-22.png", 0.55, 1.28, 8.90)
    notes(s, [
        "IMPORTANT framing - this is a probe, not a training result:",
        "- The trunk, the classifier and the text encoder are ALL frozen. Only a small Predictor trains.",
        "- So this says nothing about LoRA and nothing about our semantic training runs.",
        "- What it does say: the frozen BADAS features already encode describable scene content.",
        "THE PROCEDURE, EXACTLY:",
        "- For each held-out clip, average its window vectors and renormalise - 221 predicted vectors.",
        "- Do the same for the caption embeddings - 221 target vectors.",
        "- Both sets are L2-normalised, so the product P @ T-transpose IS the cosine similarity matrix.",
        "- For each clip take the ARGMAX along its row: the most similar caption, not the smallest loss.",
        "- A hit is when that index is the clip's own. Chance is 1 in 221 = 0.45%.",
        "- Result: 14.03%, which is 31 times chance.",
        "THE CONTROL THAT MAKES IT TRUSTWORTHY:",
        "- A predictor that ignores the video and always emits the average caption scores exactly 0.45%.",
        "- So the 31x cannot come from a shortcut - it requires actually reading the video.",
        "THE BOTTLENECK TEST (middle bar):",
        "- Rerun the identical probe on the 1,024-d pooled vector the classifier actually reads.",
        "- Still 22x chance - so caption information survives the 2560-to-1 compression.",
        "- A naive mean-pool control gives 18x, indistinguishable from 22x at n=221.",
        "- So the crash-tuned attention is not selecting against caption semantics.",
        "ON THE RIGHT PANEL:",
        "- The x-axis is training ROWS used to fit the probe: 372, 719, 1,413.",
        "- Those are the train split of the 1,761 pool; 348 rows (221 clips) are held out.",
        "- Slope 0.65 on log-log, no flattening - more captions would very likely give more."])


def s12_toolkit(prs):
    s = content_slide(prs, "Verification Toolkit - The Instruments We Built",
                      "Five measurements built to test our own explanations rather than argue "
                      "about them.", page=12)
    rows = [["Instrument", "The question it answers", "What it found"],
            ["Caption leakage gate",
             "Can a text-only classifier recover the label from the caption alone?",
             "0.964 -> 0.764 after the V12 rewrite"],
            ["Pooled-tap probe",
             "Does caption information survive the 2560x compression?",
             "Yes - still 22x chance at the bottleneck"],
            ["Delta patches vs pooled",
             "Does the semantic gradient reach the decision path, or land where pooling discards it?",
             "It reaches it - CI excludes zero"],
            ["Gradient-angle probe",
             "Do the crash and semantic objectives fight each other?",
             "No - they are near perpendicular"],
            ["Caption-similarity audit",
             "Are near-duplicate captions being punished as false negatives?",
             "Only 0.3% affected - left unchanged"]]
    table(s, rows, [1.85, 4.55, 2.60], 0.50, 1.32, 3.30, font=8.2, header_font=8.5)
    text(s, 0.50, 4.82, 9.00, 0.40,
         "The signal reaches the classifier's own representation. It simply does not help once "
         "there - so the remaining suspects are structural.", size=9.5, color=CYAN)
    notes(s, [
        "These are diagnostic tools, and each one killed a hypothesis we actually held:",
        "ON THE GRADIENT-BYPASS TEST (row 3) - where the suspicion came from:",
        "- The gradient-angle probe showed cosine near zero: the objectives were NOT fighting.",
        "- But the semantic loss still was not helping, so something had to explain the null.",
        "- Next candidate: the gradient moves directions that the attentive pooling simply discards.",
        "- Remember the pooler keeps one 1,024-d vector out of 2,560 tokens - most directions vanish.",
        "HOW WE TESTED IT:",
        "- Take the real weight difference between the control arm and a semantic arm.",
        "- Feed the same clips through both; measure the change at the patch grid and at the pooled vector.",
        "- The statistic is the ratio: how much of the change survives pooling.",
        "- Compare against a RANDOM perturbation of exactly the same magnitude, on the same clips.",
        "- Paired bootstrap over 40 clips, because real-vs-random is a paired design.",
        "THE VERDICT:",
        "- The real difference reaches the pooled vector at least as well as random does.",
        "- 95% CI on the paired difference [0.00143, 0.00163] - excludes zero.",
        "- So the gradient is NOT being routed around the decision path. That explanation is dead.",
        "ON THE SIMILARITY AUDIT (row 5) - what we did with the finding:",
        "- Concern: two clips with near-identical captions would be negatives for each other.",
        "- Measured cross-video caption cosine: mean 0.701, 99th percentile 0.870.",
        "- At a 0.90 masking threshold only about 4 of 1,413 negatives per anchor qualify - 0.3%.",
        "- We deliberately did NOT implement masking: 0.3% cannot explain a 0.02 AP gap.",
        "- Measured, then closed. It also showed the captions are genuinely clip-specific."])


def s13_future(prs):
    s = content_slide(prs, "Next Two Weeks - The Experiments That Could Clear 0.900",
                      "Ranked by expected value. All five fit inside the remaining window.",
                      page=13)
    rows = [["#", "Experiment", "What it settles", "Cost"],
            ["1", "LoRA placement:\nrouting vs content",
             "We have only ever adapted attention - WHERE the model looks. The feed-forward "
             "layers, WHAT it computes and two thirds of the weights, were never adapted.",
             "~2 h"],
            ["2", "Unfreeze the classifier,\nin both arms",
             "All five negatives ran with a classifier fitted to the ORIGINAL encoder, which "
             "LoRA then moved underneath it.", "~1 day"],
            ["3", "Redundancy measurement",
             "How much of each caption can the vision model ALREADY predict? If mostly "
             "redundant, that explains everything.", "hours"],
            ["4", "Caption the full pool\n(1,761 -> 4,446)",
             "The only lever that adds information rather than re-routing it. The scaling "
             "curve is still rising.", "~$11"],
            ["5", "Shuffled-caption control",
             "Does caption CONTENT matter, or only caption CLASS? Needed for the write-up "
             "either way.", "~1 day"]]
    table(s, rows, [0.32, 1.95, 6.05, 0.68], 0.50, 1.30, 3.45, font=7.8, header_font=8.5)
    text(s, 0.50, 4.92, 9.00, 0.35,
         "Items 1 and 2 are the realistic shots at beating 0.900. Items 3 and 5 give a "
         "defensible result whichever way they land.", size=9.5, color=GREEN)
    notes(s, [
        "Ranked by expected value; all five fit in two weeks:",
        "ITEM 1 - the newest idea, and the one I would run first:",
        "- Every run so far adapted only query/key/value - the ROUTING half of a transformer block.",
        "- The feed-forward layers hold about two thirds of the weights and were NEVER adapted.",
        "- Routing controls where the model attends; feed-forward controls what features it computes.",
        "- If the caption asks for new CONTENT but we only opened WHERE IT LOOKS, that explains a lot.",
        "- Specifically: perpendicular gradients, and a signal that reaches the classifier without helping.",
        "- Best property: it self-terminates. One epoch tells us if it is wrong, for about 12 minutes.",
        "ITEM 2 - a confound inside our own negative result:",
        "- The classifier was fitted to the original encoder and then frozen while LoRA moved it.",
        "- That penalises the semantic arm more, because it pushes features somewhere unfamiliar.",
        "- Run it in BOTH arms: if the semantic arm gains more, the signal was real but unreadable.",
        "HONEST NUMBERS IF ASKED DIRECTLY:",
        "- Beating 0.900: roughly 25-30% on item 1, roughly 45-55% on item 2.",
        "- Items 3 and 5 are not AP plays - they are what makes the write-up defensible.",
        "- After two weeks we start writing, so the cheap decisive experiments come first."])


def s14_summary(prs):
    s = content_slide(prs, "Summary - Banked and Open",
                      "What is settled, what is not, and where the next two weeks go.", page=14)
    card(s, 0.50, 1.28, 4.38, 1.72, "[BANKED]",
         ["Test AP 0.900 vs published 0.853, on 677 held-out clips",
          "Frozen features encode caption meaning at 31x chance, control at chance",
          "Six competing explanations eliminated by measurement"],
         accent=GREEN, body_size=9.5)
    card(s, 5.12, 1.28, 4.38, 1.72, "[OPEN]",
         ["Language supervision has not improved collision AP",
          "Six arms, all lost, gap widening as procedure improved",
          "The pattern points at something structural, not a bug"],
         accent=ORANGE, body_size=9.5)
    rect(s, 0.50, 3.22, 9.00, 1.62, PANEL)
    rect(s, 0.50, 3.22, 9.00, 0.05, CYAN)
    text(s, 0.68, 3.32, 8.70, 0.30, "The two leading suspects", size=12, color=CYAN, bold=True)
    text(s, 0.68, 3.68, 8.75, 1.05,
         ["- We have never adapted the parameters where the model computes content - only "
          "those that steer where it looks.",
          "- We have never let the classifier adapt to the encoder we spent every run moving.",
          "- Both are testable inside the two weeks we have, and both are cheap."],
         size=10.5, color=WHITE, space_after=5)
    notes(s, [
        "Close cleanly - keep settled and open separate:",
        "- Three results are banked and do not depend on how the language question resolves.",
        "- One thing is open, and we have been honest about it: six arms, all losing.",
        "- The gap widened each time we cleaned up execution - that points at design, not bugs.",
        "- Two specific structural suspects remain, both cheap and both testable in the window.",
        "- Question worth putting back to them: chase the positive result, or consolidate the negative?"])


def main():
    verify()
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
    for fn in (s01_title, s02_outcome, s03_idea, s04_inference, s05_training, s06_objective,
               s07_data, s08_experiments, s09_curves, s10_metrics, s11_probe, s12_toolkit,
               s13_future, s14_summary):
        fn(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"  wrote {OUT}")
    print(f"  {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    main()
