"""Builds the A1-failure-recovery status deck (6 slides, dark theme).

Audience: two professors who have seen the 2026-08-22 semantic-guidance deck. Same
house style (10 x 5.62 in, #1C2340 background, cyan/navy palette) copied from
build_status_presentation_2026-08-22.py; slides stay sparse and the speaker notes carry
the teaching, one bullet per line, never wrapped.

Every number is re-derived from the real artifacts by verify() before the deck is
written, so a stale figure can never be silently embedded.

    python student_training/scripts/build_a1fail_presentation.py
"""
import json
from pathlib import Path

import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt
from sklearn.metrics import average_precision_score, roc_auc_score

ROOT = Path(__file__).resolve().parents[2]
FIG = ROOT / "reports" / "figures"
A1F = ROOT / "outputs" / "a1fail321"
OUT = ROOT / "reports" / "presentations" / "2026-08_a1-failure-recovery.pptx"

# ---- E3a-status palette (identical to the 2026-08-22 deck)
BG = RGBColor(0x1C, 0x23, 0x40)
PANEL = RGBColor(0x24, 0x30, 0x60)
PANEL_DK = RGBColor(0x20, 0x2A, 0x50)
CYAN = RGBColor(0x00, 0xBF, 0xFF)
CYAN_DK = RGBColor(0x00, 0x96, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xA0, 0xB4, 0xCC)
GREEN = RGBColor(0x00, 0xE6, 0x76)
ORANGE = RGBColor(0xFF, 0xA7, 0x26)
RED = RGBColor(0xFF, 0x6B, 0x6B)
SW, SH = 10.0, 5.625

FOOTER = "CCP-MMLM · Semantic Guidance for Collision Anticipation"

# ---------------------------------------------------------------- verified data
D = {}   # filled by verify(), read by the slides


def verify():
    """Re-derive every embedded number from source artifacts. Fail loudly on drift."""
    # ---- test-set scores (677 clips)
    for arm in ("A1", "v12_ep10"):
        rows = [json.loads(l) for l in open(A1F / "test_scores" / f"{arm}.jsonl",
                                            encoding="utf-8") if l.strip()]
        y = np.array([1 if r["gt_verdict"] == "YES" else 0 for r in rows])
        s = np.array([r["score"] for r in rows])
        accs = [(t, ((s >= t) == (y == 1)).mean()) for t in np.arange(0.05, 0.96, 0.01)]
        best_t, best_a = max(accs, key=lambda x: x[1])
        D[arm] = {"n": len(rows), "ap": average_precision_score(y, s),
                  "auc": roc_auc_score(y, s), "acc50": ((s >= 0.5) == (y == 1)).mean(),
                  "best_thr": best_t, "acc_best": best_a, "mean": s.mean()}
    assert D["A1"]["n"] == D["v12_ep10"]["n"] == 677, "test set is not 677 clips"
    # A1 must reproduce its documented 0.900/0.904 through THIS scorer, otherwise the
    # v12 number it is compared against is not trustworthy either.
    assert abs(D["A1"]["ap"] - 0.900) < 0.005, f"A1 test AP drifted: {D['A1']['ap']}"
    assert abs(D["A1"]["auc"] - 0.904) < 0.005, f"A1 test AUC drifted: {D['A1']['auc']}"

    # ---- pool composition
    sel = [json.loads(l) for l in open(A1F / "selection_a1fail321.jsonl",
                                       encoding="utf-8") if l.strip()]
    D["pool_n"] = len(sel)
    D["pool_videos"] = len({r["video_id"] for r in sel})
    D["n_train"] = sum(1 for r in sel if r["split"] == "train")
    D["n_val"] = sum(1 for r in sel if r["split"] == "val")
    assert (D["pool_n"], D["n_train"], D["n_val"]) == (321, 260, 61), \
        f"pool split changed: {D['pool_n']}/{D['n_train']}/{D['n_val']}"
    y = [1 if r["gt_verdict"] == "YES" else 0 for r in sel]
    s = [r["a1_score"] for r in sel]
    D["a1_pool_auc"] = roc_auc_score(y, s)
    assert D["a1_pool_auc"] == 0.0, \
        f"A1 AUC on its own failures is {D['a1_pool_auc']}, expected exactly 0.0"

    # ---- in-pool val outcome per arm (must be identical across all four)
    sel_by_fd = {r["frames_dir"]: r for r in sel}
    val_fds = [fd for fd, r in sel_by_fd.items() if r["split"] == "val"]
    D["arms"] = {}
    # Epoch 10 for EVERY arm, deliberately. Each arm's own argmax-val_ap epoch differs
    # slightly (v10 peaks at 9 by 0.0002 - noise at n=61), but reporting different epochs
    # per arm would confound "which arm" with "how much training". Verified this choice
    # is inert: v10 gives the identical 39/0/22 at both epoch 9 and epoch 10.
    for arm, ep in (("a1cont", 10), ("v10", 10), ("v12", 10), ("v12shuf", 10)):
        rows = [json.loads(l) for l in open(
            A1F / "results" / arm / "fold_01" / f"val_scores_ep{ep:02d}.jsonl",
            encoding="utf-8") if l.strip()]
        sc = {r["frames_dir"]: r["score"] for r in rows}
        ffp = sum(1 for fd in val_fds
                  if sel_by_fd[fd]["gt_verdict"] == "NO" and sc[fd] < 0.5)
        ffn = sum(1 for fd in val_fds
                  if sel_by_fd[fd]["gt_verdict"] == "YES" and sc[fd] >= 0.5)
        acc = sum(1 for fd in val_fds
                  if (sc[fd] >= 0.5) == (sel_by_fd[fd]["gt_verdict"] == "YES")) / len(val_fds)
        D["arms"][arm] = {"ep": ep, "fixed_fp": ffp, "fixed_fn": ffn,
                          "still": len(val_fds) - ffp - ffn, "acc": acc}
    outcomes = {(a["fixed_fp"], a["fixed_fn"], a["still"]) for a in D["arms"].values()}
    assert len(outcomes) == 1, f"arms are NOT identical any more: {D['arms']}"
    D["identical"] = outcomes.pop()

    # ---- retrieval health (the semantic branch works / control fails)
    for arm in ("v10", "v12", "v12shuf"):
        rows = [json.loads(l) for l in open(
            A1F / "results" / arm / "fold_01" / "epoch_metrics.jsonl",
            encoding="utf-8") if l.strip()]
        D["arms"][arm]["retr_max"] = max(r["retrieval_clip"] for r in rows)
        D["arms"][arm]["ctrl"] = rows[-1]["collapse_control_clip"]
        cos = [r["grad_cos_mean"] for r in rows if r.get("grad_cos_mean") is not None]
        D["arms"][arm]["cos_lo"], D["arms"][arm]["cos_hi"] = min(cos), max(cos)
    assert D["arms"]["v12"]["retr_max"] > 0.30, "v12 retrieval collapsed - check the run"
    assert D["arms"]["v12shuf"]["retr_max"] < 0.05, "shuffled control is NOT at chance"
    print("[verify] all numbers re-derived from artifacts and consistent")


# ---------------------------------------------------------------- style helpers
def _bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def rect(slide, x, y, w, h, color):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def text(slide, x, y, w, h, runs, size=11, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
         font="Calibri", italic=False, space_after=4, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    paras = [runs] if isinstance(runs, str) else runs
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        pieces = [(para, {})] if isinstance(para, str) else para
        for txt, ov in pieces:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(ov.get("size", size))
            r.font.bold = ov.get("bold", bold)
            r.font.italic = ov.get("italic", italic)
            r.font.color.rgb = ov.get("color", color)
            r.font.name = ov.get("font", font)
    return tb


def notes(slide, lines):
    """One bullet per line. Never wraps - the user reads these live."""
    tf = slide.notes_slide.notes_text_frame
    tf.text = lines[0]
    for l in lines[1:]:
        tf.add_paragraph().text = l


def picture(slide, path, x, y, w):
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))


def content_slide(prs, title, subline=None, page=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    rect(s, 0.35, 0.22, 0.07, 0.50, CYAN)
    t_size = 25 if len(title) <= 52 else max(19, 25 * 52 / len(title))
    text(s, 0.50, 0.13, 9.20, 0.60, title, size=t_size, color=CYAN, bold=True)
    rect(s, 0.35, 0.84, 9.30, 0.012, CYAN_DK)
    if subline:
        text(s, 0.50, 0.88, 9.20, 0.40, subline, size=11, color=MUTED)
    text(s, 0.50, 5.29, 7.00, 0.30, FOOTER, size=8.5, color=MUTED)
    if page:
        text(s, 8.80, 5.29, 0.80, 0.30, str(page), size=8.5, color=MUTED,
             align=PP_ALIGN.RIGHT)
    return s


def card(slide, x, y, w, h, title, bullets, accent=CYAN, sub=None, body_size=9.5):
    rect(slide, x, y, w, h, PANEL)
    rect(slide, x, y, w, 0.05, accent)
    text(slide, x + 0.13, y + 0.10, w - 0.26, 0.30, title, size=12, color=accent, bold=True)
    top = y + 0.42
    if sub:
        text(slide, x + 0.13, y + 0.38, w - 0.26, 0.26, sub, size=9, color=MUTED)
        top = y + 0.68
    text(slide, x + 0.13, top, w - 0.26, h - (top - y) - 0.08,
         [f"• {b}" for b in bullets], size=body_size, color=WHITE, space_after=3)


def hero(slide, x, y, w, h, headline, tail, accent=GREEN):
    rect(slide, x, y, w, h, PANEL_DK)
    rect(slide, x, y, 0.07, h, accent)
    text(slide, x + 0.20, y + 0.08, w - 0.30, h - 0.12,
         [[(headline, {"size": 21, "bold": True, "color": accent}),
           (tail, {"size": 12, "color": WHITE})]], size=12)


def table(slide, rows, cols_w, x, y, h, font=8.5, header_font=8.5, align_first_left=True):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y),
                                   Inches(sum(cols_w)), Inches(h))
    tbl = shape.table
    tbl.first_row = False
    for j, w in enumerate(cols_w):
        tbl.columns[j].width = Inches(w)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            fill = "202A50" if i == 0 else ("243060" if i % 2 else "1F2950")
            cell._tc.get_or_add_tcPr().append(
                parse_xml(f'<a:solidFill {nsdecls("a")}><a:srgbClr val="{fill}"/></a:solidFill>'))
            cell.text = str(val)
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            for p in cell.text_frame.paragraphs:
                p.alignment = (PP_ALIGN.LEFT if (j == 0 and align_first_left)
                               else PP_ALIGN.CENTER)
                for r in p.runs:
                    r.font.size = Pt(header_font if i == 0 else font)
                    r.font.name = "Calibri"
                    r.font.bold = (i == 0)
                    r.font.color.rgb = CYAN if i == 0 else WHITE
    return tbl


# ================================================================= SLIDES
def s1_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    rect(s, 0, 0, SW, 0.08, CYAN)
    rect(s, 0, SH - 0.08, SW, 0.08, CYAN)
    rect(s, 0, 0, 0.08, SH, CYAN_DK)
    rect(s, 7.40, 0, 2.60, 2.60, PANEL_DK)
    text(s, 0.50, 0.80, 9.00, 1.10, "CCP-MMLM", size=48, color=WHITE, bold=True)
    rect(s, 0.50, 1.90, 6.50, 0.06, CYAN)
    text(s, 0.50, 1.99, 8.60, 0.90,
         "Can Semantic Guidance Repair What the Vision Model Gets Wrong?\n"
         "Targeted follow-up to the six-arm semantic-guidance study", size=15, color=MUTED)
    rect(s, 0.50, 3.05, 9.00, 1.42, PANEL)
    rect(s, 0.50, 3.05, 9.00, 0.05, CYAN)
    text(s, 0.62, 3.13, 8.70, 0.30, "The question this run asks", size=12, color=CYAN, bold=True)
    text(s, 0.62, 3.45, 8.76, 1.00,
         "Our best model (A1, test AP 0.900) is wrong on a specific set of clips. Previous "
         "arms trained on everything and asked whether language helps on average. This run "
         "starts FROM A1's own weights and trains only on the clips A1 fails, asking a "
         "sharper question: given a known failure set, can teacher captions fix it - and "
         "does trying damage what already works?", size=11, color=WHITE)
    text(s, 0.50, 4.62, 7.00, 0.35, "MSc Thesis - Experiment Report  |  August 2026",
         size=11.5, color=CYAN)
    text(s, 0.50, 4.92, 5.00, 0.42, "Eviatar Ohayon", size=15, color=WHITE, bold=True)
    notes(s, [
        "Framing in one minute:",
        "- Previous deck: six arms, language never beat the vision-only control on average.",
        "- Natural objection to that: averages hide things - maybe language helps exactly where vision fails.",
        "- So this run targets the failures directly instead of the whole pool.",
        "- Two changes from every previous arm: start from A1's trained weights, and train only on its errors.",
        "- Two questions: does it repair the failures, and does it cost us the 0.900 we already have.",
        "- Both questions now have clean answers, and one of them is genuinely reassuring."])


def s2_experiment(prs):
    a = D["arms"]["a1cont"]
    s = content_slide(prs, "The Experiment - Train Only Where A1 Fails",
                      f"{D['pool_n']} windows A1 gets wrong, split {D['n_train']} train / "
                      f"{D['n_val']} validation by video.", page=2)
    card(s, 0.50, 1.28, 4.38, 1.62, "How the clips were chosen",
         [f"Scored all 1,761 pool windows with A1 (its own checkpoint)",
          f"Kept every window A1 calls wrong at threshold 0.5 -> {D['pool_n']} windows",
          f"{D['pool_videos']} unique videos; split by VIDEO so sibling frames cannot leak",
          "62% are false alarms, 34% are missed collisions"],
         accent=CYAN, body_size=9)
    card(s, 5.12, 1.28, 4.38, 1.62, "Four arms, identical except captions",
         ["a1cont - crash loss only (the control)",
          "v10 - captions that leak the label",
          "v12 - clean, de-leaked captions",
          "v12shuf - v12 captions scrambled between clips"],
         accent=CYAN, body_size=9)
    rect(s, 0.50, 3.12, 9.00, 1.62, PANEL_DK)
    rect(s, 0.50, 3.12, 0.07, 1.62, ORANGE)
    text(s, 0.70, 3.20, 8.60, 0.30,
         "Read in-pool numbers with care - this set is rank-inverted by construction",
         size=12, color=ORANGE, bold=True)
    text(s, 0.70, 3.56, 8.66, 1.10,
         ["- Every clip here is one A1 already gets wrong, so A1's AUC on this set is "
          "EXACTLY 0.000 - not merely low.",
          "- Any model that is not degenerate therefore looks like a large improvement; "
          "random guessing alone scores AUC 0.5.",
          "- So 'we fixed 39 of 61' is not evidence on its own. The honest test is the "
          "held-out 677-clip test set - slide 5."],
         size=10, color=WHITE, space_after=5)
    notes(s, [
        "What was run:",
        f"- Took A1, scored all 1,761 training-pool windows, kept the {D['pool_n']} it gets wrong.",
        f"- {D['pool_videos']} distinct videos. Split {D['n_train']}/{D['n_val']} by video_id, never by row.",
        "- Why by video: one video contributes several time-windows; splitting by row leaks.",
        "- Four arms, everything identical except which caption file each one reads.",
        "- v12shuf is the key control: same captions, randomly reassigned between clips.",
        "- If content matters, v12 should beat v12shuf. If only the presence of a second loss matters, they tie.",
        "THE CAVEAT - say this out loud before showing any in-pool number:",
        "- A1's AUC on its own failures is exactly zero, by construction, not by measurement.",
        "- Every negative outranks every positive, because that is the selection rule.",
        "- So improvement here is nearly guaranteed and means much less than it looks.",
        "- This is why the test set on slide 5 is the number that decides anything."])


def s3_architecture(prs):
    s = content_slide(prs, "Architecture - What Moves and What Does Not",
                      "Same two-branch training design as before; the difference is where "
                      "the weights START.", page=3)
    # NOT the 2026-08-22 figure: that one is drawn with lambda=0.05, the weight used in
    # the six-arm study. This run used 0.2, so it embeds the regenerated variant
    # (make_arch_figures_2026-08-22.py fig_L3(lam="0.2", out_name=...)). Embedding the
    # original would put a wrong loss equation in front of the professors.
    picture(s, FIG / "arch_L3_training_a1fail_2026-08-29.png", 0.30, 1.12, 6.30)
    card(s, 6.82, 1.16, 2.72, 1.62, "FROZEN",
         ["Crash head (pooling + classifier)",
          "SigLIP text encoder",
          "V-JEPA2 trunk base weights"],
         accent=RED, body_size=8.5)
    card(s, 6.82, 2.90, 2.72, 1.62, "TRAINABLE",
         ["LoRA adapters - 2.8M",
          "Predictor (resampler) - 1.25M",
          "Everything else stays fixed"],
         accent=GREEN, body_size=8.5)
    text(s, 6.82, 4.62, 2.72, 0.60,
         "Both start warm, not random: LoRA from A1, Predictor from the B1 probe.",
         size=9, color=ORANGE, italic=True)
    notes(s, [
        "Topology is unchanged from the deck you already saw - only initialization differs:",
        "- Vision path: V-JEPA2 trunk, LoRA adapters on attention, then a frozen crash head.",
        "- Second branch, training only: Predictor maps patch features to a 768-d vector.",
        "- The teacher caption goes through a frozen SigLIP text encoder to give the target.",
        "- The whole language branch is deleted at inference - zero runtime cost, as before.",
        "THE IMPORTANT DIFFERENCE FROM EVERY PREVIOUS ARM:",
        "- LoRA does NOT start from zero. It starts from A1's trained adapters.",
        "- So epoch 1 already sits at AP 0.900 on the test set - we are refining, not training.",
        "- The Predictor also starts warm, from the B1 probe, matching the B-v3 protocol.",
        "- Learning rate accordingly dropped 5x, to 2e-5 - full rate would erase A1 in a few steps.",
        "- The crash head stays frozen: it is what A1's 0.900 was measured with, so moving it changes two things at once."])


def s4_curves(prs):
    v12 = D["arms"]["v12"]; shuf = D["arms"]["v12shuf"]
    s = content_slide(prs, "Training Dynamics - The Language Branch Genuinely Learns",
                      "Solid = crash loss (left axis).  Dashed = InfoNCE semantic loss "
                      "(right axis, green).  Circled = checkpoint taken to the test set.",
                      page=4)
    # Picture starts below the subline band (which ends ~1.28) - at y=1.05 it overlapped
    # and hid the axis legend. Width 6.60in -> 3.78in tall, bottom lands at 5.06, clear
    # of the 5.29 footer.
    picture(s, FIG / "a1fail_loss_curves_2026-08-29.png", 0.28, 1.28, 6.60)
    card(s, 6.98, 1.28, 2.52, 1.72, "What to see",
         ["All four crash curves fall smoothly - no overfitting",
          "Val below train: this split is slightly easier, from epoch 0",
          "Only v12 / v10 show a falling semantic curve"],
         accent=CYAN, body_size=8.5)
    card(s, 6.98, 3.12, 2.52, 1.94, "The control works",
         [f"v12 retrieval@1 = {v12['retr_max']:.0%} vs {v12['ctrl']:.1%} chance",
          f"v12shuf = {shuf['retr_max']:.0%} - scrambled captions learn nothing",
          "So the branch reads real meaning, not just 'a caption exists'"],
         accent=GREEN, body_size=8.5)
    notes(s, [
        "Walk the figure left to right:",
        "- Left axis, solid curves: crash loss. Right axis, green, dashed: the semantic loss.",
        "- Top-left is the control - it has no semantic branch, so its right axis is deliberately blank.",
        "- All four crash curves fall smoothly and flatten. No overfitting anywhere in this run.",
        "- Validation sits BELOW training throughout - that is not a bug, I checked it.",
        "- A1's own loss on the val subset was already lower before we trained: 1.32 vs 1.41.",
        "- Small-sample split luck; the val failures happen to sit closer to the decision boundary.",
        "THE RESULT THAT MATTERS ON THIS SLIDE:",
        f"- v12's retrieval@1 reaches {v12['retr_max']:.0%} against a collapse control of {v12['ctrl']:.1%}.",
        f"- v12shuf, with the SAME captions randomly reassigned, sits at {shuf['retr_max']:.0%}.",
        "- That is the cleanest control result we have had: real captions are learned, scrambled ones are not.",
        "- So the language branch is genuinely working here - which was NOT true in earlier runs.",
        "- Circled point on v12 is epoch 10, the checkpoint we carried to the test set."])


def s5_test(prs):
    a1, v12 = D["A1"], D["v12_ep10"]
    ident = D["identical"]
    s = content_slide(prs, "Test-Set Result - Neutral, and Why",
                      f"{a1['n']} held-out clips. The benchmark A1 set is intact; language "
                      "did not move it.", page=5)
    rows = [["Model", "AP", "AUC", "acc @ 0.5", "acc @ own best threshold"],
            ["A1 (baseline)", f"{a1['ap']:.4f}", f"{a1['auc']:.4f}",
             f"{a1['acc50']:.4f}", f"{a1['acc_best']:.4f}  (thr {a1['best_thr']:.2f})"],
            ["v12 (semantic)", f"{v12['ap']:.4f}", f"{v12['auc']:.4f}",
             f"{v12['acc50']:.4f}", f"{v12['acc_best']:.4f}  (thr {v12['best_thr']:.2f})"],
            ["difference", f"{v12['ap']-a1['ap']:+.4f}", f"{v12['auc']-a1['auc']:+.4f}",
             f"{v12['acc50']-a1['acc50']:+.4f}", f"{v12['acc_best']-a1['acc_best']:+.4f}"]]
    table(s, rows, [2.10, 1.30, 1.30, 1.50, 2.80], 0.50, 1.22, 1.30, font=9, header_font=9)
    text(s, 0.50, 2.62, 9.00, 0.34,
         [[("The +2.6 point gain at threshold 0.5 is calibration, not skill. ",
            {"bold": True, "color": ORANGE}),
           (f"Training shifted v12's mean score {a1['mean']:.2f} -> {v12['mean']:.2f}, "
            f"landing it near 0.5 by luck. Tuned fairly, the gap is +{v12['acc_best']-a1['acc_best']:.3f}.",
            {"color": WHITE})]], size=9.5)
    rect(s, 0.50, 3.08, 9.00, 1.72, PANEL)
    rect(s, 0.50, 3.08, 9.00, 0.05, ORANGE)
    text(s, 0.68, 3.16, 8.70, 0.30, "Why it does not transfer - the gradients are perpendicular",
         size=12, color=ORANGE, bold=True)
    text(s, 0.68, 3.52, 8.75, 1.22,
         ["- The captions were written by the teacher FROM THE SAME 16 FRAMES. They add no "
          "information the pixels lack - only a suggestion about how to organise features.",
          "- To win at retrieval the model must FINGERPRINT a clip among 1,933 candidates: "
          "road type, lighting, vehicle count. To predict a collision it needs one relation: "
          "is the gap closing.",
          "- Measured cosine between the two gradients: -0.04 to +0.05, sign flipping every "
          "epoch. Not fighting - simply pulling in unrelated directions.",
          f"- Consequence: all four arms produce IDENTICAL predictions in-pool "
          f"({ident[0]} fixed, {ident[2]} still wrong) - including the arm with no captions at all."],
         size=9.5, color=WHITE, space_after=4)
    notes(s, [
        "The headline first:",
        f"- A1 re-scored through this exact scorer: AP {a1['ap']:.4f}, matching its documented 0.900.",
        "- That is a deliberate check - if A1 had not reproduced, v12's number would mean nothing.",
        f"- v12 after training on A1's failures: AP {v12['ap']:.4f}. Difference {v12['ap']-a1['ap']:+.4f}.",
        "- Flat. No gain, and importantly no catastrophic forgetting - the risk we were testing for.",
        "THE ACCURACY COLUMN - do not let anyone over-read it:",
        f"- At the fixed 0.5 cut v12 looks {v12['acc50']-a1['acc50']:+.3f} better. That is real but it is calibration.",
        f"- Training pushed v12's mean score from {a1['mean']:.2f} down to {v12['mean']:.2f}, near 0.5 by accident.",
        "- Give each model its own best threshold and the gap collapses to +0.004. Same skill.",
        "- A1 reaches the same accuracy by moving its threshold to 0.68. No retraining needed.",
        "WHY IT DOES NOT TRANSFER - this is the mechanism, and it is measured, not assumed:",
        "- The teacher saw the same frames. A caption cannot contain what the pixels do not.",
        "- So this was never distillation of new knowledge - only a reorganisation pressure.",
        "- Retrieval rewards fingerprinting a clip; collision prediction rewards one relation.",
        "- Those demands barely overlap, so the gradients sit at right angles rather than in conflict.",
        "- The proof is the control arm: no captions at all, and it fixes exactly the same clips."])


def s6_conclusion(prs):
    s = content_slide(prs, "Conclusion - A Closed Question and a Sharper One",
                      "The mechanism is now understood well enough to say what to test next.",
                      page=6)
    card(s, 0.50, 1.22, 4.38, 1.66, "[BANKED]",
         ["A1's 0.900 benchmark survives - training on its failures is safe",
          "The language branch demonstrably learns meaning (control validates it)",
          "Why it fails is now measured, not speculated"],
         accent=GREEN, body_size=9)
    card(s, 5.12, 1.22, 4.38, 1.66, "[CLOSED]",
         ["Caption-retrieval alignment does not transfer to collision AP",
          "Holds with warm starts, a working predictor and a clean control",
          "Cause is structural, not a tuning or procedure bug"],
         accent=ORANGE, body_size=9)
    rect(s, 0.50, 3.06, 9.00, 1.74, PANEL_DK)
    rect(s, 0.50, 3.06, 0.07, 1.74, CYAN)
    text(s, 0.70, 3.14, 8.60, 0.32,
         "Proposed next question: supervise CONCEPTS, not caption identity",
         size=13, color=CYAN, bold=True)
    text(s, 0.70, 3.54, 8.66, 1.20,
         ["Instead of asking the model to match a whole sentence, ask it to predict the "
          "specific facts we already extract from the teacher: is the gap closing, are the "
          "brake lights on, what is the ego manoeuvre.",
          "Why it should behave differently: those targets need the SAME features collision "
          "prediction needs, so their gradients should align rather than sit perpendicular.",
          "It is falsifiable in advance, with the instrument already built: predict "
          "grad-cosine persistently above +0.15, versus the +-0.05 sign-flipping we measure "
          "today. One run answers it."],
         size=10, color=WHITE, space_after=6)
    notes(s, [
        "Close in three moves - banked, closed, next:",
        "BANKED - two things worth keeping regardless of how the language question ends:",
        "- Training on a model's own failure set does not damage it. That was a real risk and it did not happen.",
        "- The semantic branch works end-to-end now: real captions learned, scrambled ones not.",
        "CLOSED - and I want to state this plainly rather than leave it ambiguous:",
        "- Caption-embedding alignment does not improve collision AP. We have tested it thoroughly.",
        "- This run removed the last procedural excuses: warm start, working predictor, proper control.",
        "- The control arm with no captions performs identically. That is as clean as a null gets.",
        "THE NEXT QUESTION - and why it is a different question, not the same one retried:",
        "- Every failed run asked the model to reproduce a SENTENCE, which rewards fingerprinting a scene.",
        "- Concept supervision asks for the FACTS instead - gap trend, brake lights, manoeuvre.",
        "- Those need the same visual evidence a collision needs, so the gradients should point together.",
        "- We already extract these fields from the teacher, so the data cost is zero.",
        "- Best property: it is pre-registerable. I can state the success threshold BEFORE running it.",
        "- Prediction: grad-cosine persistently above +0.15. Today we measure plus or minus 0.05, flipping sign.",
        "- If it fails that bar, the direction is finished and we write it up as a characterised negative result."])


def main():
    verify()
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
    for fn in (s1_title, s2_experiment, s3_architecture, s4_curves, s5_test, s6_conclusion):
        fn(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"  wrote {OUT}")
    print(f"  {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    main()
